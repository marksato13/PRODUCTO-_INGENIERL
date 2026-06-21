"""
generar_slides_defensa.py
Slides de defensa PPI — UPeU 2026
Sistema de Deteccion Temprana de Comportamientos Anomalos en Redes de Datos
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree
import copy, os

# ── Paleta ────────────────────────────────────────────────────────────────────
AZUL     = RGBColor(0x00, 0x4A, 0x99)   # UPeU azul institucional
AZUL_OSC = RGBColor(0x00, 0x2D, 0x62)
AZUL_CLR = RGBColor(0xD6, 0xE4, 0xF5)
VERDE    = RGBColor(0x1A, 0x7C, 0x4A)
ROJO     = RGBColor(0xC0, 0x20, 0x20)
GRIS     = RGBColor(0x55, 0x55, 0x55)
GRIS_CLR = RGBColor(0xF2, 0xF4, 0xF8)
BLANCO   = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO    = RGBColor(0x1A, 0x1A, 0x1A)
AMARILLO = RGBColor(0xF5, 0xA6, 0x23)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completamente en blanco

# ── Helpers ───────────────────────────────────────────────────────────────────
def slide():
    return prs.slides.add_slide(BLANK)

def rect(sl, x, y, w, h, fill=None, border=None):
    shp = sl.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.line.fill.background()
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp

def txbox(sl, text, x, y, w, h, size=18, bold=False, color=NEGRO,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def header_bar(sl, title, subtitle=None):
    """Barra azul oscuro arriba."""
    r = rect(sl, 0, 0, W, Inches(1.4), fill=AZUL_OSC)
    txbox(sl, title, Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.75),
          size=28, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(sl, subtitle, Inches(0.35), Inches(0.85), Inches(12.5), Inches(0.45),
              size=14, color=AZUL_CLR, align=PP_ALIGN.LEFT)

def footer(sl, num, total=17):
    rect(sl, 0, H - Inches(0.28), W, Inches(0.28), fill=AZUL)
    txbox(sl, f"PPI — UPeU 2026  |  Rubén Salazar Tocas  |  {num}/{total}",
          Inches(0.3), H - Inches(0.28), Inches(12.5), Inches(0.28),
          size=9, color=BLANCO, align=PP_ALIGN.CENTER)

def bullet(sl, items, x, y, w, size=15, color=NEGRO, spacing=0.38):
    """Lista de bullets simples."""
    cy = y
    for item in items:
        prefix = "• "
        indent = 0
        if item.startswith("  "):
            prefix = "   –"
            indent = Inches(0.2)
            item = item.strip()
        txbox(sl, prefix + item, x + indent, cy, w - indent, Inches(spacing + 0.05),
              size=size, color=color)
        cy += Inches(spacing)

def kpi_box(sl, label, value, x, y, w=Inches(2.1), h=Inches(1.2),
            bg=AZUL, fg=BLANCO, sub=None):
    rect(sl, x, y, w, h, fill=bg)
    txbox(sl, value, x, y + Inches(0.05), w, Inches(0.7),
          size=30, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txbox(sl, label, x, y + Inches(0.7), w, Inches(0.35),
          size=11, color=fg, align=PP_ALIGN.CENTER)
    if sub:
        txbox(sl, sub, x, y + Inches(0.95), w, Inches(0.22),
              size=9, color=fg, align=PP_ALIGN.CENTER)

def tabla(sl, headers, rows, x, y, col_ws, row_h=Inches(0.38),
          hdr_bg=AZUL_OSC, hdr_fg=BLANCO, alt=AZUL_CLR):
    cx = x
    # headers
    for i, (h, cw) in enumerate(zip(headers, col_ws)):
        rect(sl, cx, y, cw, row_h, fill=hdr_bg)
        txbox(sl, h, cx + Inches(0.05), y + Inches(0.05), cw - Inches(0.1), row_h,
              size=11, bold=True, color=hdr_fg, align=PP_ALIGN.CENTER)
        cx += cw
    # rows
    for ri, row in enumerate(rows):
        cy = y + row_h * (ri + 1)
        cx = x
        bg = alt if ri % 2 == 0 else BLANCO
        for ci, (cell, cw) in enumerate(zip(row, col_ws)):
            rect(sl, cx, cy, cw, row_h, fill=bg, border=AZUL_CLR)
            txbox(sl, cell, cx + Inches(0.05), cy + Inches(0.04), cw - Inches(0.1), row_h,
                  size=11, color=NEGRO, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
            cx += cw

def divider(sl, y, color=AZUL_CLR):
    rect(sl, Inches(0.35), y, W - Inches(0.7), Pt(1.5), fill=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=AZUL_OSC)                            # fondo azul
rect(sl, 0, H - Inches(2.0), W, Inches(2.0), fill=AZUL)        # banda inferior

# Banda de color decorativa
rect(sl, 0, Inches(1.4), W, Inches(0.08), fill=AMARILLO)

# Institución
txbox(sl, "UNIVERSIDAD PERUANA UNIÓN", 0, Inches(0.18), W, Inches(0.55),
      size=14, bold=True, color=AZUL_CLR, align=PP_ALIGN.CENTER)
txbox(sl, "Facultad de Ingeniería y Arquitectura  ·  PPI 2026",
      0, Inches(0.68), W, Inches(0.45),
      size=12, color=AZUL_CLR, align=PP_ALIGN.CENTER)

# Título principal
txbox(sl,
      "Sistema de Detección Temprana de\nComportamientos Anómalos en Redes de Datos",
      Inches(0.6), Inches(1.7), W - Inches(1.2), Inches(1.9),
      size=34, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

txbox(sl, "Mediante Isolation Forest + Control Inline con iptables/ipset",
      Inches(0.6), Inches(3.5), W - Inches(1.2), Inches(0.55),
      size=17, color=AZUL_CLR, align=PP_ALIGN.CENTER, italic=True)

# Datos
txbox(sl, "Autor:     Rubén Mark Salazar Tocas",
      Inches(3.2), Inches(4.5), Inches(7), Inches(0.45), size=14, color=BLANCO)
txbox(sl, "Asesor 1:  Ing. Nemias Saboya Rios",
      Inches(3.2), Inches(4.9), Inches(7), Inches(0.45), size=14, color=BLANCO)
txbox(sl, "Asesor 2:  Ing. Fernando Manuel Asin Gomez",
      Inches(3.2), Inches(5.3), Inches(7), Inches(0.45), size=14, color=BLANCO)
txbox(sl, "Junio 2026",
      Inches(3.2), Inches(5.85), Inches(7), Inches(0.45), size=14, color=AMARILLO, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Agenda", "Estructura de la presentación")
footer(sl, 2)

items_izq = [
    "1.  Problema y motivación",
    "2.  Objetivos del proyecto",
    "3.  Marco teórico",
    "4.  Topología del laboratorio",
    "5.  Arquitectura del sistema",
    "6.  Recolección y procesamiento de datos",
    "7.  Modelo Isolation Forest",
    "8.  Umbrales de decisión",
]
items_der = [
    "9.  Motor de decisión en tiempo real",
    "10. Validación F6 — 40 corridas",
    "11. Resultados vs. requisitos",
    "12. Dashboard de monitoreo",
    "13. Módulo Predictor — XGBoost",
    "14. Conclusiones",
    "15. Trabajo futuro",
    "16. Preguntas",
]

bullet(sl, items_izq, Inches(0.5), Inches(1.6), Inches(6.0), size=15, spacing=0.37)
bullet(sl, items_der, Inches(6.9), Inches(1.6), Inches(6.0), size=15, spacing=0.37)
divider(sl, Inches(1.55))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problema y Motivación
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Problema y Motivación", "¿Por qué es necesario este sistema?")
footer(sl, 3)

# Panel izquierdo
rect(sl, Inches(0.3), Inches(1.55), Inches(6.0), Inches(5.65), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "Contexto universitario", Inches(0.4), Inches(1.65), Inches(5.8), Inches(0.4),
      size=13, bold=True, color=AZUL_OSC)
bullet(sl, [
    "Redes abiertas con cientos de dispositivos",
    "Tráfico legítimo mezclado con ataques",
    "Sin mecanismo de detección automática",
    "Respuesta manual lenta (minutos-horas)",
], Inches(0.4), Inches(2.1), Inches(5.6), size=13, spacing=0.42)

txbox(sl, "Tipos de ataque presentes", Inches(0.4), Inches(3.9), Inches(5.8), Inches(0.4),
      size=13, bold=True, color=AZUL_OSC)
bullet(sl, [
    "SYN Flood — satura el servidor",
    "Port Scan — reconocimiento de servicios",
    "UDP/ICMP Flood — agota el ancho de banda",
    "Brute Force SSH — compromete credenciales",
    "HTTP Abuse — deniega el servicio web",
], Inches(0.4), Inches(4.35), Inches(5.6), size=13, spacing=0.35)

# Panel derecho — estadísticas
rect(sl, Inches(6.7), Inches(1.55), Inches(6.3), Inches(5.65), fill=AZUL_OSC)
txbox(sl, "El problema en números", Inches(6.85), Inches(1.7), Inches(6.0), Inches(0.45),
      size=14, bold=True, color=BLANCO)

stats = [
    ("~62s", "Lead Time sin detección automática"),
    ("100%", "Disponibilidad requerida del servicio"),
    ("0%", "Interrupción por Tiempo de Límite (ITL)"),
    ("<500ms", "Latencia máxima permitida por flujo"),
]
cy = Inches(2.3)
for val, lbl in stats:
    rect(sl, Inches(7.0), cy, Inches(5.5), Inches(0.95), fill=AZUL)
    txbox(sl, val, Inches(7.1), cy + Inches(0.05), Inches(2.0), Inches(0.5),
          size=24, bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)
    txbox(sl, lbl, Inches(9.2), cy + Inches(0.2), Inches(3.1), Inches(0.6),
          size=12, color=BLANCO)
    cy += Inches(1.15)

txbox(sl, "Propuesta: detección automática en tiempo real\ncon machine learning + control inline",
      Inches(6.85), Inches(6.7), Inches(6.0), Inches(0.55),
      size=12, color=AMARILLO, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Objetivos
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Objetivos", "General y específicos")
footer(sl, 4)

rect(sl, Inches(0.3), Inches(1.55), W - Inches(0.6), Inches(1.05), fill=AZUL_OSC)
txbox(sl, "OBJETIVO GENERAL", Inches(0.5), Inches(1.6), Inches(4.0), Inches(0.4),
      size=11, bold=True, color=AMARILLO)
txbox(sl, "Implementar y validar un sistema de detección temprana de comportamientos anómalos en redes de datos"
      " universitarias utilizando Isolation Forest, con control inline mediante iptables/ipset.",
      Inches(0.5), Inches(1.92), W - Inches(1.0), Inches(0.55),
      size=13, color=BLANCO)

txbox(sl, "OBJETIVOS ESPECÍFICOS", Inches(0.5), Inches(2.85), Inches(6.0), Inches(0.4),
      size=11, bold=True, color=AZUL_OSC)
divider(sl, Inches(3.2))

oes = [
    ("F1–F2", "Capturar y preprocesar tráfico real (normal + anómalo) con Suricata 7.0.3. Extraer 14 features por flujo."),
    ("F3",    "Entrenar Isolation Forest (n=300) y derivar umbrales τ1/τ2 óptimos mediante curva ROC."),
    ("F4–F5", "Implementar motor de decisión en tiempo real con latencia < 500 ms y control inline (ipset)."),
    ("F6",    "Validar el sistema con 40 corridas: Disponibilidad ≥ 99%, ITL = 0%, AUC-ROC ≥ 0.85."),
]
cy = Inches(3.3)
for fase, texto in oes:
    rect(sl, Inches(0.4), cy, Inches(0.95), Inches(0.72), fill=AZUL)
    txbox(sl, fase, Inches(0.4), cy + Inches(0.14), Inches(0.95), Inches(0.45),
          size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txbox(sl, texto, Inches(1.5), cy + Inches(0.1), Inches(11.3), Inches(0.55),
          size=13, color=NEGRO)
    cy += Inches(0.85)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Marco Teórico
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Marco Teórico", "Tecnologías y algoritmos clave")
footer(sl, 5)

conceptos = [
    ("Isolation Forest (IF)",
     "Algoritmo de detección de anomalías basado en árboles de aislamiento. "
     "Asigna scores de anomalía a cada flujo: score cercano a −1 = muy anómalo, "
     "score cercano a 0 = normal. Entrenado solo con datos normales (one-class)."),
    ("Suricata 7.0.3",
     "Motor IDS/IPS open source. Captura flujos de red en formato eve.json "
     "con metadatos por flujo: bytes, paquetes, puertos, protocolo, duración. "
     "Opera en modo pasivo sobre la interfaz ens35 del sensor."),
    ("ipset / iptables",
     "Mecanismo de control inline en el servidor (192.168.0.120). "
     "ppi_blocked (DROP) y ppi_limited (hashlimit 100pkt/s). "
     "Operaciones atómicas O(1) — no afectan latencia del motor."),
    ("Youden Index (J = TPR − FPR)",
     "Criterio para seleccionar el umbral óptimo τ1. Maximiza la diferencia "
     "entre sensibilidad (TPR) y tasa de falsos positivos (FPR). "
     "Resultado: τ1 = −0.4459, TPR = 99.40%, FPR = 20.47%."),
]
cy = Inches(1.62)
for titulo, texto in conceptos:
    rect(sl, Inches(0.3), cy, W - Inches(0.6), Inches(1.25), fill=BLANCO, border=AZUL_CLR)
    txbox(sl, titulo, Inches(0.5), cy + Inches(0.07), Inches(12.0), Inches(0.38),
          size=13, bold=True, color=AZUL_OSC)
    txbox(sl, texto, Inches(0.5), cy + Inches(0.42), Inches(12.3), Inches(0.72),
          size=12, color=GRIS)
    cy += Inches(1.35)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Topología del Laboratorio
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Topología del Laboratorio", "5 máquinas virtuales — red 192.168.0.0/24")
footer(sl, 6)

# Tabla de nodos
tabla(sl,
    ["IP", "VM / Rol", "Función"],
    [
        ["192.168.0.10",  "Win 11         — Cliente",       "Genera tráfico de usuario final"],
        ["192.168.0.20",  "Ubuntu Desktop — Administrador", "Origen tráfico normal (Grupo A), Claude Code"],
        ["192.168.0.100", "Kali Linux     — Atacante",      "Origen tráfico anómalo (Grupo B)"],
        ["192.168.0.110", "Ubuntu Sensor  — IDS",           "Suricata 7.0.3 en ens35, motor de decisión"],
        ["192.168.0.120", "Ubuntu Server  — Servicio",      "nginx:80, SSH:22, ipset ppi_blocked/limited"],
    ],
    Inches(0.35), Inches(1.65),
    [Inches(1.7), Inches(3.5), Inches(7.6)],
    row_h=Inches(0.5)
)

# Diagrama de flujo simple
cy = Inches(4.45)
nodos = [
    (Inches(0.4),  "Kali\n.100",    ROJO),
    (Inches(3.0),  "Desktop\n.20",  VERDE),
    (Inches(5.8),  "Sensor\n.110",  AZUL),
    (Inches(8.8),  "Motor IF\n(sensor)", AZUL_OSC),
    (Inches(11.5), "Server\n.120",  VERDE),
]
bw = Inches(1.7)
bh = Inches(0.9)
for x, lbl, col in nodos:
    rect(sl, x, cy, bw, bh, fill=col)
    txbox(sl, lbl, x, cy + Inches(0.12), bw, bh,
          size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# Flechas (texto simulado)
for i in range(len(nodos) - 1):
    xa = nodos[i][0] + bw
    xb = nodos[i+1][0]
    xm = xa + (xb - xa) / 2 - Inches(0.3)
    rect(sl, xa, cy + Inches(0.4), xb - xa, Pt(2), fill=AZUL_CLR)
    txbox(sl, "→", xm, cy + Inches(0.25), Inches(0.6), Inches(0.35),
          size=14, color=AZUL, align=PP_ALIGN.CENTER)

txbox(sl, "Flujo de datos: tráfico → Suricata (eve.json) → Motor IF → decisión → ipset (server)",
      Inches(0.35), Inches(5.65), W - Inches(0.7), Inches(0.4),
      size=11, color=GRIS, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Arquitectura / Pipeline
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Arquitectura del Sistema", "Pipeline de 6 fases (F1–F6)")
footer(sl, 7)

fases = [
    ("F1", "Captura",        "Suricata\neve.json",          AZUL),
    ("F2", "Captura",        "3 grupos\n.gz directo",        AZUL),
    ("F3", "Modelado",       "IsoForest\nAUC=0.8998",       AZUL_OSC),
    ("F4", "Motor",          "motor_decision\n.py",          AZUL_OSC),
    ("F5", "Control",        "ipset\nBLOCK/LIMIT",          ROJO),
    ("F6", "Validación",     "40 corridas\nDisp=100%",      VERDE),
]
bw = Inches(1.85)
bh = Inches(1.7)
gap = Inches(0.3)
cx  = Inches(0.35)
cy  = Inches(1.85)

for fase, titulo, desc, color in fases:
    rect(sl, cx, cy, bw, bh, fill=color)
    txbox(sl, fase, cx, cy + Inches(0.1), bw, Inches(0.45),
          size=20, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txbox(sl, titulo, cx, cy + Inches(0.5), bw, Inches(0.45),
          size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txbox(sl, desc, cx, cy + Inches(0.95), bw, Inches(0.7),
          size=11, color=AZUL_CLR, align=PP_ALIGN.CENTER)
    if fase != "F6":
        txbox(sl, "→", cx + bw, cy + Inches(0.65), gap + Inches(0.05), Inches(0.5),
              size=20, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    cx += bw + gap

# Descripción abajo
rect(sl, Inches(0.35), Inches(3.8), W - Inches(0.7), Inches(3.45), fill=BLANCO, border=AZUL_CLR)

desc_items = [
    ("F1–F2", "Suricata captura flows en eve.json. Captura separada: Grupo A=normal puro (Kali apagada), B=ataques puros (Desktop quieto), C=mixto (motor detenido). Los .gz son la fuente directa para F3 — sin CSVs intermedios."),
    ("F3",    "Isolation Forest (n=300, contamination=0.05) entrenado solo con flujos normales. "
              "fase3_evaluar.py evalúa holdout 20% + Grupo B → ROC → τ1/τ2 → metricas_offline.txt (fuente única)."),
    ("F4–F5", "motor_decision.py hace tail de eve.json en tiempo real. Por cada flujo: extrae features → "
              "score IF → compara τ1/τ2 → envía a enforce.sh → SSH al servidor → ipset add/del."),
    ("F6",    "f6_corridas.py ejecuta 40 corridas en 4 grupos (Normal, Mixto, Reeval, Final). "
              "Mide Disponibilidad, ITL, Lead Time, Latencia."),
]
cy2 = Inches(3.95)
for fase, texto in desc_items:
    txbox(sl, f"[{fase}]  {texto}", Inches(0.55), cy2, Inches(12.3), Inches(0.72),
          size=12, color=NEGRO)
    cy2 += Inches(0.73)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Datos y Features
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Recolección y Procesamiento de Datos", "F1–F2: Captura → 14 features por flujo")
footer(sl, 8)

# Escenarios
rect(sl, Inches(0.3), Inches(1.6), Inches(4.5), Inches(5.55), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "Escenarios de tráfico", Inches(0.4), Inches(1.7), Inches(4.3), Inches(0.4),
      size=13, bold=True, color=AZUL_OSC)
bullet(sl, [
    "Grupo A — Normal (desde Desktop .20)",
    "  A1 HTTP normal (curl/wget → :80)",
    "  A2 SSH legítimo (8 min)",
    "  A3 Transferencia scp/wget",
    "  A4 Tráfico sostenido mixto",
    "Grupo B — Anómalo (desde Kali .100)",
    "  B1 SYN Flood (hping3)",
    "  B2 Port Scan (nmap -sS)",
    "  B3 UDP Flood (hping3 --udp)",
    "  B4 ICMP Flood",
    "  B5 HTTP Abuse (curl bucle)",
    "  B6 Brute Force SSH (hydra)",
    "Grupo C — Mixto (Desktop + Kali)",
    "  C1/C2/C3 simultáneos",
], Inches(0.4), Inches(2.15), Inches(4.1), size=12, spacing=0.33)

# Features y estadísticas
rect(sl, Inches(5.1), Inches(1.6), Inches(7.9), Inches(2.55), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "14 Features del modelo", Inches(5.25), Inches(1.7), Inches(7.5), Inches(0.4),
      size=13, bold=True, color=AZUL_OSC)
feat_text = (
    "pkts_toserver, pkts_toclient, bytes_toserver, bytes_toclient,\n"
    "duration, pkt_rate, byte_rate, pkt_ratio, byte_ratio,\n"
    "avg_pkt_size, is_tcp, is_udp, is_icmp, dest_port"
)
txbox(sl, feat_text, Inches(5.25), Inches(2.1), Inches(7.5), Inches(0.85),
      size=12, color=NEGRO)
txbox(sl, "Fuente: eve.json (Suricata) → campos app_proto, proto, flow, pkts, bytes, ts",
      Inches(5.25), Inches(2.95), Inches(7.5), Inches(0.45),
      size=11, color=GRIS, italic=True)

# Estadísticas del dataset
rect(sl, Inches(5.1), Inches(4.35), Inches(7.9), Inches(2.8), fill=AZUL_OSC)
txbox(sl, "Dataset final", Inches(5.25), Inches(4.45), Inches(7.5), Inches(0.4),
      size=13, bold=True, color=BLANCO)
stats2 = [
    ("67,135", "Flujos normales   (28 capturas Jun-02, 04, 15)"),
    ("598,285","Flujos anómalos   (13 capturas Jun-02, 04, 15)"),
    ("80/20",   "Split entrenamiento/holdout (aleatorio, seed=42)"),
    ("53,708", "Flujos usados en entrenamiento del IF"),
]
cy3 = Inches(4.95)
for val, lbl in stats2:
    txbox(sl, val, Inches(5.3), cy3, Inches(1.5), Inches(0.38),
          size=14, bold=True, color=AMARILLO)
    txbox(sl, lbl, Inches(6.85), cy3 + Inches(0.05), Inches(5.9), Inches(0.38),
          size=12, color=BLANCO)
    cy3 += Inches(0.52)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Modelo Isolation Forest
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Modelo Isolation Forest", "F3: Entrenamiento y evaluación — AUC-ROC = 0.8998")
footer(sl, 9)

# KPIs del modelo
kpis_modelo = [
    ("AUC-ROC",    "0.8998",  AZUL_OSC, BLANCO),
    ("Precisión",  "99.54%",  VERDE,    BLANCO),
    ("Recall",     "99.40%",  VERDE,    BLANCO),
    ("F1-Score",   "0.9947",  VERDE,    BLANCO),
    ("n_estim.",   "300",     AZUL,     BLANCO),
]
cx = Inches(0.35)
for lbl, val, bg, fg in kpis_modelo:
    kpi_box(sl, lbl, val, cx, Inches(1.62), w=Inches(2.4), h=Inches(1.2), bg=bg, fg=fg)
    cx += Inches(2.55)

# Parámetros del modelo
rect(sl, Inches(0.35), Inches(3.05), Inches(5.4), Inches(4.15), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "Hiperparámetros", Inches(0.5), Inches(3.15), Inches(5.0), Inches(0.4),
      size=13, bold=True, color=AZUL_OSC)
tabla(sl,
    ["Parámetro", "Valor"],
    [
        ["n_estimators",   "300"],
        ["contamination",  "0.05"],
        ["random_state",   "42"],
        ["max_samples",    "'auto'"],
        ["sklearn",        "1.9.0"],
    ],
    Inches(0.5), Inches(3.6),
    [Inches(2.5), Inches(2.5)],
    row_h=Inches(0.44)
)

# Curva ROC y scores
rect(sl, Inches(6.0), Inches(3.05), Inches(7.0), Inches(4.15), fill=AZUL_OSC)
txbox(sl, "Separación de scores por clase", Inches(6.15), Inches(3.15), Inches(6.7), Inches(0.4),
      size=13, bold=True, color=BLANCO)
scores = [
    ("Flujos normales",  "−0.3965 ± 0.0753", "Score alto (cerca de 0)"),
    ("Flujos anómalos", "−0.5420 ± 0.0900", "Score bajo (cerca de −1)"),
    ("Δ separación",     "0.1454",            "Diferencia media entre clases"),
]
cy4 = Inches(3.65)
for lbl, val, nota in scores:
    txbox(sl, lbl, Inches(6.2), cy4, Inches(2.8), Inches(0.35), size=11, color=AZUL_CLR)
    txbox(sl, val, Inches(8.9), cy4, Inches(2.2), Inches(0.35), size=13, bold=True, color=AMARILLO)
    txbox(sl, nota, Inches(6.2), cy4 + Inches(0.32), Inches(4.7), Inches(0.32),
          size=10, color=AZUL_CLR, italic=True)
    cy4 += Inches(0.9)

txbox(sl, "El modelo fue re-entrenado el 2026-06-16 con todos los datos disponibles\n"
          "(28+13 archivos) corrigiendo bug de fecha hardcodeada — AUC mejoró de 0.8955 a 0.8998.",
      Inches(0.35), Inches(6.8), W - Inches(0.7), Inches(0.5),
      size=11, color=GRIS, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Umbrales de Decisión
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Umbrales de Decisión τ1 / τ2", "Derivados de la curva ROC — criterios estadísticos justificados")
footer(sl, 10)

# Tabla principal de umbrales
tabla(sl,
    ["Umbral", "Valor", "Criterio", "TPR", "FPR", "Acción"],
    [
        ["τ1 (PERMIT/LIMIT)", "−0.4459", "Youden J = TPR−FPR máximo", "99.40%", "20.47%", "score > τ1 → PERMIT"],
        ["τ2 (LIMIT/BLOCK)",  "−0.6027", "FPR ≤ 2% con TPR máximo",  "18.27%",  "1.99%", "τ2 < score ≤ τ1 → LIMIT"],
        ["— (BLOCK)",         "≤ τ2",    "Anomalía confirmada",        "—",       "—",    "score ≤ τ2 → BLOCK (DROP)"],
    ],
    Inches(0.35), Inches(1.65),
    [Inches(2.0), Inches(1.3), Inches(3.0), Inches(1.05), Inches(1.05), Inches(4.7)],
    row_h=Inches(0.55)
)

# Panel: lógica de decisión
rect(sl, Inches(0.35), Inches(3.5), Inches(5.8), Inches(3.65), fill=AZUL_OSC)
txbox(sl, "Lógica del motor por flujo", Inches(0.5), Inches(3.6), Inches(5.5), Inches(0.4),
      size=13, bold=True, color=BLANCO)
logica = [
    "1.  IP en whitelist? → IGNORAR",
    "2.  score > −0.4459 (τ1)? → PERMIT",
    "3.  score > −0.6027 (τ2)? → LIMIT",
    "          hashlimit 100 pkt/s",
    "          ipset ppi_limited",
    "4.  score ≤ −0.6027 → BLOCK",
    "          DROP iptables",
    "          ipset ppi_blocked",
]
cy5 = Inches(4.1)
for item in logica:
    col = AMARILLO if item.strip().startswith(("1.", "2.", "3.", "4.")) else AZUL_CLR
    sz  = 13 if col == AMARILLO else 11
    txbox(sl, item, Inches(0.5), cy5, Inches(5.4), Inches(0.38),
          size=sz, color=col)
    cy5 += Inches(0.38)

# Panel: nota sobre FPR
rect(sl, Inches(6.5), Inches(3.5), Inches(6.5), Inches(3.65), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "¿Por qué FPR = 20.47% es aceptable?", Inches(6.65), Inches(3.6), Inches(6.2), Inches(0.4),
      size=13, bold=True, color=AZUL_OSC)
nota_text = (
    "Reducir FPR a ≤ 5% requeriría τ1 = −0.5547.\n\n"
    "Problema: SYN floods obtienen score ≈ −0.49.\n"
    "Con τ1 = −0.555, los SYN floods serían PERMIT\n"
    "(no detectados).\n\n"
    "El Youden Index (τ1 = −0.4459) preserva\n"
    "la detección de SYN Flood al costo de\n"
    "falsos positivos en tráfico muy anómalo-looking.\n\n"
    "Mitigación: whitelist de IPs conocidas\n"
    "(Desktop .20, Sensor .110, Server .120)."
)
txbox(sl, nota_text, Inches(6.65), Inches(4.1), Inches(6.1), Inches(2.9),
      size=12, color=NEGRO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Motor de Decisión en Tiempo Real
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Motor de Decisión en Tiempo Real", "F4–F5: Latencia P95 = 34.8 ms  (<500 ms requerido)")
footer(sl, 11)

# Diagrama de pipeline (cajas horizontales)
pasos = [
    ("tail\neve.json", AZUL),
    ("parser\nfeatures", AZUL),
    ("IF\nscore", AZUL_OSC),
    ("τ1/τ2\ndecision", AZUL_OSC),
    ("enforce\n.sh", ROJO),
    ("SSH\nserver", ROJO),
    ("ipset\nBLOCK", ROJO),
]
bw2  = Inches(1.6)
bh2  = Inches(1.1)
gap2 = Inches(0.2)
cx2  = Inches(0.3)
cy2  = Inches(1.75)
for lbl, col in pasos:
    rect(sl, cx2, cy2, bw2, bh2, fill=col)
    txbox(sl, lbl, cx2, cy2 + Inches(0.15), bw2, bh2,
          size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    if lbl != pasos[-1][0]:
        txbox(sl, "→", cx2 + bw2, cy2 + Inches(0.35), gap2 + Inches(0.1), Inches(0.4),
              size=14, bold=True, color=AZUL, align=PP_ALIGN.CENTER)
    cx2 += bw2 + gap2

# Latencia por etapa (aproximada)
lats = ["~1ms", "~5ms", "~25ms", "~1ms", "~1ms", "~30ms SSH", "~1ms"]
cx3  = Inches(0.3)
for lat in lats:
    txbox(sl, lat, cx3, cy2 + bh2 + Inches(0.05), bw2, Inches(0.35),
          size=10, color=GRIS, align=PP_ALIGN.CENTER, italic=True)
    cx3 += bw2 + gap2

# Detectores heurísticos
rect(sl, Inches(0.3), Inches(3.4), Inches(6.3), Inches(3.75), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "Detectores Heurísticos (complemento al IF)", Inches(0.45), Inches(3.5),
      Inches(6.0), Inches(0.4), size=13, bold=True, color=AZUL_OSC)
bullet(sl, [
    "SSH Brute Force",
    "  ≥ 5 intentos / 60s → LIMIT",
    "  ≥ 15 intentos / 60s → BLOCK",
    "HTTP Abuse",
    "  ≥ 50 req / 30s → LIMIT",
    "  ≥ 100 req / 30s → BLOCK",
    "Ventana deslizante en memoria (deque)",
    "Opera EN PARALELO con el IF",
], Inches(0.5), Inches(4.0), Inches(5.9), size=12, spacing=0.36)

# Métricas de rendimiento
rect(sl, Inches(6.9), Inches(3.4), Inches(6.1), Inches(3.75), fill=AZUL_OSC)
txbox(sl, "Métricas de rendimiento", Inches(7.05), Inches(3.5), Inches(5.8), Inches(0.4),
      size=13, bold=True, color=BLANCO)
met = [
    ("Latencia media",  "34.5 ms"),
    ("Latencia P95",    "34.8 ms"),
    ("Umbral req.",     "< 500 ms"),
    ("Veredicto",       "CUMPLE"),
    ("Lead Time detec.","~62 s"),
    ("Disponibilidad",  "100%"),
]
cy6 = Inches(4.0)
for lbl, val in met:
    col_val = AMARILLO if val in ("CUMPLE", "100%") else AZUL_CLR
    txbox(sl, lbl, Inches(7.05), cy6, Inches(3.0), Inches(0.38), size=12, color=AZUL_CLR)
    txbox(sl, val, Inches(9.9), cy6, Inches(2.8), Inches(0.38), size=13, bold=True, color=col_val)
    cy6 += Inches(0.48)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Validación F6
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Validación F6 — 40 Corridas", "Metodología de validación experimental exhaustiva")
footer(sl, 12)

# KPIs principales
kpis_f6 = [
    ("Disponibilidad", "100%",    VERDE,    BLANCO, "corridas 1-40"),
    ("ITL Global",     "0%",      VERDE,    BLANCO, "sin interrupciones"),
    ("Lead Time",      "61.92s",  AZUL_OSC, BLANCO, "SYN Flood C11"),
    ("Latencia P95",   "34.8ms",  AZUL,     BLANCO, "req. <500ms"),
    ("AUC-ROC",        "0.8998",  AZUL,     BLANCO, "req. ≥0.85"),
]
cx7 = Inches(0.3)
for lbl, val, bg, fg, sub in kpis_f6:
    kpi_box(sl, lbl, val, cx7, Inches(1.62), w=Inches(2.4), h=Inches(1.3), bg=bg, fg=fg, sub=sub)
    cx7 += Inches(2.55)

# Grupos de corridas
tabla(sl,
    ["Grupo", "Corridas", "Descripción", "Resultado"],
    [
        ["Normal",   "1–10",  "Solo tráfico legítimo (Desktop → Server)",      "Disponibilidad = 100%"],
        ["Mixto",    "11–20", "Tráfico normal + ataque simultáneo (C11=SYN)",   "Detección en 61.92s  "],
        ["Reeval",   "21–30", "Re-evaluación post-bloqueo de 192.168.0.100",    "ITL = 0%, Kali bloq. "],
        ["Final",    "31–40", "Escenario completo sistema en producción",        "Todas métricas OK    "],
    ],
    Inches(0.35), Inches(3.1),
    [Inches(1.3), Inches(1.4), Inches(5.9), Inches(4.6)],
    row_h=Inches(0.5)
)

txbox(sl, "Nota metodológica: corridas 12-40 muestran flows_anom=0 porque 192.168.0.100 permanece en el set "
          "in-memory 'bloqueados' tras la detección en corrida 11. Comportamiento CORRECTO — el atacante "
          "queda contenido sin reactivarse.",
      Inches(0.35), Inches(5.38), W - Inches(0.7), Inches(0.65),
      size=11, color=GRIS, italic=True)

# Mini gráfica de texto
rect(sl, Inches(0.35), Inches(6.15), W - Inches(0.7), Inches(1.05), fill=AZUL_OSC)
txbox(sl, "Timeline de detección (corrida 11 — SYN Flood):",
      Inches(0.5), Inches(6.22), Inches(4.5), Inches(0.38), size=11, bold=True, color=BLANCO)
txbox(sl, "t=0s  Motor iniciado  |  t≈60s  Suricata cierra flow TCP half-open  |"
          "  t=61.92s  IF score=−0.71 ≤ τ2  →  BLOCK aplicado vía ipset",
      Inches(0.5), Inches(6.57), W - Inches(0.9), Inches(0.5),
      size=12, color=AMARILLO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Resultados vs. Requisitos
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Resultados vs. Requisitos", "Todos los criterios de aceptación CUMPLIDOS")
footer(sl, 13)

tabla(sl,
    ["Métrica", "Requisito", "Resultado", "Estado"],
    [
        ["AUC-ROC",                   "≥ 0.85",       "0.8998",    "CUMPLE ✓"],
        ["Precisión",                 "≥ 95%",        "99.54%",    "CUMPLE ✓"],
        ["Recall",                    "≥ 95%",        "99.40%",    "CUMPLE ✓"],
        ["F1-Score",                  "≥ 0.90",       "0.9947",    "CUMPLE ✓"],
        ["Latencia P95 por flujo",    "< 500 ms",     "34.8 ms",   "CUMPLE ✓"],
        ["Disponibilidad del servicio","≥ 99%",        "100%",      "CUMPLE ✓"],
        ["ITL (Interr. por Tiempo)",  "= 0%",         "0%",        "CUMPLE ✓"],
        ["Lead Time de detección",    "< 120 s",      "61.92 s",   "CUMPLE ✓"],
        ["Corridas exitosas F6",      "40/40",        "40/40",     "CUMPLE ✓"],
    ],
    Inches(0.35), Inches(1.65),
    [Inches(4.2), Inches(2.0), Inches(2.0), Inches(4.8)],
    row_h=Inches(0.48)
)

txbox(sl, "TODOS LOS REQUISITOS CUMPLIDOS — El sistema es operacionalmente viable.",
      Inches(0.35), Inches(6.82), W - Inches(0.7), Inches(0.45),
      size=14, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Dashboard Web
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Dashboard de Monitoreo", "Flask + SSE — Tiempo real en puerto 8080")
footer(sl, 14)

# Panel izquierdo — descripción
rect(sl, Inches(0.3), Inches(1.6), Inches(5.2), Inches(5.6), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "Características del dashboard", Inches(0.45), Inches(1.7), Inches(5.0), Inches(0.4),
      size=13, bold=True, color=AZUL_OSC)
bullet(sl, [
    "Flask + Server-Sent Events (SSE)",
    "Actualización push sin polling",
    "Endpoints REST:",
    "  /api/stats    → métricas globales",
    "  /api/stream   → eventos SSE",
    "  /api/alerts   → últimas 100 alertas",
    "  /api/block    → bloqueo manual",
    "  /api/unblock  → desbloqueo manual",
    "Sin memory leaks (validado >22h)",
    "Acceso: http://192.168.0.110:8080",
], Inches(0.45), Inches(2.2), Inches(4.9), size=12, spacing=0.37)

# Panel derecho — mockup del dashboard
rect(sl, Inches(5.8), Inches(1.6), Inches(7.2), Inches(5.6), fill=NEGRO)
txbox(sl, "◉  Dashboard PPI — Detección Anomalías",
      Inches(5.9), Inches(1.65), Inches(7.0), Inches(0.4),
      size=11, bold=True, color=VERDE)
rect(sl, Inches(5.9), Inches(2.1), Inches(7.0), Inches(0.9), fill=AZUL_OSC)
kpi_data = [("Flows", "124,892"), ("Alerts", "3"), ("Blocked", "1"), ("Uptime", "22h")]
cx8 = Inches(5.95)
for lbl, val in kpi_data:
    txbox(sl, val, cx8, Inches(2.12), Inches(1.6), Inches(0.48),
          size=16, bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)
    txbox(sl, lbl, cx8, Inches(2.6), Inches(1.6), Inches(0.3),
          size=9, color=AZUL_CLR, align=PP_ALIGN.CENTER)
    cx8 += Inches(1.72)

rect(sl, Inches(5.9), Inches(3.1), Inches(7.0), Inches(0.35), fill=AZUL_OSC)
txbox(sl, "Alertas recientes", Inches(5.95), Inches(3.12), Inches(6.9), Inches(0.3),
      size=10, bold=True, color=BLANCO)
alerts = [
    "[19:38:42] BLOCK  192.168.0.100  score=-0.71  SYN Flood",
    "[19:38:01] LIMIT  192.168.0.100  score=-0.53  sospechoso",
    "[17:22:15] PERMIT 192.168.0.20   score=-0.31  normal",
]
cy9 = Inches(3.5)
for a in alerts:
    col = ROJO if "BLOCK" in a[0] else (AMARILLO if "LIMIT" in a[0] else VERDE)
    txbox(sl, a[0], Inches(5.9), cy9, Inches(7.1), Inches(0.38), size=10, color=col)
    cy9 += Inches(0.4)

rect(sl, Inches(5.9), Inches(4.75), Inches(7.0), Inches(2.1), fill=RGBColor(0x11,0x1A,0x2E))
txbox(sl, "Histograma de scores (últimas 5min)",
      Inches(5.95), Inches(4.8), Inches(6.9), Inches(0.35), size=10, color=AZUL_CLR)
bars = [0.4, 0.7, 0.9, 0.6, 0.3, 0.8, 1.0, 0.5, 0.4, 0.6]
for i, h in enumerate(bars):
    bh3 = Inches(h * 1.35)
    by  = Inches(6.8) - bh3
    col = ROJO if h > 0.75 else (AMARILLO if h > 0.5 else VERDE)
    rect(sl, Inches(6.0) + Inches(i * 0.63), by, Inches(0.5), bh3, fill=col)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Módulo Predictor Temporal — XGBoost
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Módulo Predictor Temporal — XGBoost",
           "Anticipación predictiva integrada al motor IF (proceso paralelo)")
footer(sl, 15)

# Panel izquierdo — cómo funciona
rect(sl, Inches(0.3), Inches(1.55), Inches(6.0), Inches(5.65), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "Señal predictiva — gap entre estadísticas",
      Inches(0.45), Inches(1.65), Inches(5.7), Inches(0.38), size=12, bold=True, color=AZUL_OSC)
bullet(sl, [
    "Motor escribe stats cada 500 flows",
    "Normal:  gap ~ 174s   (tráfico bajo)",
    "Ataque:  gap ~  17s   (10× más rápido)",
    "Predictor mide cuánto tarda la próxima stats",
    "→ detecta aceleración ANTES del BLOCK",
], Inches(0.5), Inches(2.1), Inches(5.6), size=12, spacing=0.38)

txbox(sl, "Modelo seleccionado — XGBoost (AUC=0.58)",
      Inches(0.45), Inches(3.55), Inches(5.7), Inches(0.38), size=12, bold=True, color=AZUL_OSC)
bullet(sl, [
    "ARIMA:        AUC=0.50  (descartado)",
    "Random Forest: AUC=0.48  (descartado)",
    "XGBoost:      AUC=0.58  ✓ ELEGIDO",
    "scale_pos_weight compensa desbalance",
    "P ≥ 0.70  →  ALERTA-PREDICTIVA (log + SSE)",
], Inches(0.5), Inches(4.0), Inches(5.6), size=12, spacing=0.38)

# Panel derecho — resultados corridas
rect(sl, Inches(6.5), Inches(1.55), Inches(6.5), Inches(3.3), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "Corridas de validación (P4-P10)",
      Inches(6.65), Inches(1.65), Inches(6.2), Inches(0.38), size=12, bold=True, color=AZUL_OSC)
tabla(sl,
    ["Corrida", "Tipo", "Resultado", "P%"],
    [("P4", "Ataque", "ALERTA ✓", "81.43%"),
     ("P5", "Ataque", "ALERTA ✓", "83.46%"),
     ("P7", "Normal", "Sin alerta ✓", "—"),
     ("P8", "Normal", "Sin alerta ✓", "—"),
     ("P9-P10", "Normal", "Sin alerta ✓", "—")],
    Inches(6.55), Inches(2.1),
    [Inches(1.4), Inches(1.4), Inches(2.2), Inches(1.2)], row_h=Inches(0.36))

# KPIs
rect(sl, Inches(6.5), Inches(5.0), Inches(6.5), Inches(2.2), fill=AZUL_OSC)
txbox(sl, "Métricas del predictor", Inches(6.65), Inches(5.05), Inches(6.2), Inches(0.35),
      size=12, bold=True, color=BLANCO)
kpi_pred = [("TPR ataques", "2/2 = 100%"), ("FPR normal", "0/4 = 0%"),
            ("P media ataque", "82.4%"), ("Servicios", "3 activos")]
cy_k = Inches(5.45)
for lbl, val in kpi_pred:
    txbox(sl, f"  {lbl}:", Inches(6.55), cy_k, Inches(3.3), Inches(0.35),
          size=12, color=AZUL_CLR)
    txbox(sl, val, Inches(9.85), cy_k, Inches(3.0), Inches(0.35),
          size=12, bold=True, color=AMARILLO)
    cy_k += Inches(0.42)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Conclusiones
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Conclusiones", "Lo que se logró — impacto y contribución")
footer(sl, 16)

conclusiones = [
    ("Sistema funcional end-to-end",
     "Se implementó y validó un pipeline completo: captura (Suricata) → preproceso → "
     "Isolation Forest → motor tiempo real → control inline (ipset). "
     "Todas las fases F1-F6 operativas en laboratorio real."),
    ("Métricas de calidad superiores",
     "AUC-ROC = 0.8998 (req. ≥ 0.85), Precisión = 99.54%, Recall = 99.40%, F1 = 0.9947. "
     "El modelo detecta anomalías con alta confianza sobre 14 features de capa de red."),
    ("Respuesta automática en tiempo real",
     "Latencia P95 = 34.8ms (req. <500ms, mejora de 14×). Lead Time de detección ≈ 62s "
     "limitado por el timeout de Suricata para flows TCP half-open."),
    ("Disponibilidad garantizada",
     "40/40 corridas con Disponibilidad = 100% e ITL = 0%. "
     "El mecanismo de whitelist previene auto-bloqueo del sistema."),
    ("Defensa justificada del umbral FPR=20.47%",
     "Reducir FPR a ≤5% sacrificaría la detección de SYN Flood (score≈−0.49 > τ1=−0.555). "
     "El índice de Youden es el umbral metodológicamente correcto en este contexto."),
]
cy10 = Inches(1.65)
for i, (titulo, texto) in enumerate(conclusiones):
    rect(sl, Inches(0.3), cy10, Inches(0.5), Inches(0.95), fill=VERDE)
    txbox(sl, str(i+1), Inches(0.3), cy10 + Inches(0.18), Inches(0.5), Inches(0.6),
          size=20, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    rect(sl, Inches(0.9), cy10, W - Inches(1.2), Inches(0.95), fill=BLANCO, border=AZUL_CLR)
    txbox(sl, titulo, Inches(1.05), cy10 + Inches(0.06), Inches(11.8), Inches(0.35),
          size=12, bold=True, color=AZUL_OSC)
    txbox(sl, texto, Inches(1.05), cy10 + Inches(0.4), Inches(11.8), Inches(0.48),
          size=11, color=NEGRO)
    cy10 += Inches(1.08)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Trabajo Futuro y Preguntas
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, W, H, fill=GRIS_CLR)
header_bar(sl, "Trabajo Futuro  /  Preguntas", "Líneas de extensión del sistema")
footer(sl, 17)

rect(sl, Inches(0.3), Inches(1.6), Inches(6.3), Inches(5.6), fill=BLANCO, border=AZUL_CLR)
txbox(sl, "Trabajo Futuro", Inches(0.45), Inches(1.7), Inches(6.0), Inches(0.4),
      size=14, bold=True, color=AZUL_OSC)
futuro = [
    "Ensemble IF + AE (AND gate): FPR -49%, F1 +4.8pp",
    "Re-entrenamiento incremental online",
    "Extender whitelist dinámica (DHCP/LDAP)",
    "Soporte multi-sensor (agente distribuido)",
    "Integración con SIEM (Wazuh / Elastic)",
    "Notificaciones automáticas (email/Slack)",
    "Predictor con ventana de series temporales (LSTM)",
    "Despliegue en contenedor Docker",
]
bullet(sl, futuro, Inches(0.5), Inches(2.2), Inches(5.9), size=13, spacing=0.41)

# Panel preguntas
rect(sl, Inches(6.9), Inches(1.6), Inches(6.1), Inches(5.6), fill=AZUL_OSC)
txbox(sl, "¿Preguntas?", Inches(7.05), Inches(2.0), Inches(5.7), Inches(0.7),
      size=28, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
txbox(sl, "Rubén Mark Salazar Tocas",
      Inches(7.05), Inches(3.1), Inches(5.7), Inches(0.45),
      size=14, color=AZUL_CLR, align=PP_ALIGN.CENTER)
txbox(sl, "makosdfrs@gmail.com",
      Inches(7.05), Inches(3.55), Inches(5.7), Inches(0.4),
      size=13, color=AMARILLO, align=PP_ALIGN.CENTER)

txbox(sl, "Repositorio:",
      Inches(7.05), Inches(4.3), Inches(5.7), Inches(0.35),
      size=11, color=AZUL_CLR, align=PP_ALIGN.CENTER)
txbox(sl, "github.com/marksato13/PRODUCTO-_INGENIERL",
      Inches(7.05), Inches(4.65), Inches(5.7), Inches(0.35),
      size=11, color=BLANCO, align=PP_ALIGN.CENTER)

txbox(sl, "Motor en línea:",
      Inches(7.05), Inches(5.35), Inches(5.7), Inches(0.35),
      size=11, color=AZUL_CLR, align=PP_ALIGN.CENTER)
txbox(sl, "http://192.168.0.110:8080",
      Inches(7.05), Inches(5.7), Inches(5.7), Inches(0.35),
      size=11, color=BLANCO, align=PP_ALIGN.CENTER)
txbox(sl, "(Dashboard en vivo durante la defensa)",
      Inches(7.05), Inches(6.05), Inches(5.7), Inches(0.35),
      size=10, color=AMARILLO, italic=True, align=PP_ALIGN.CENTER)

# ── Guardar ───────────────────────────────────────────────────────────────────
OUT = "/home/m4rk/ppi-surikata-producto/results/slides_defensa_PPI_UPeU_2026.pptx"
prs.save(OUT)
print(f"OK: {OUT}")
import os
print(f"Tamano: {os.path.getsize(OUT)//1024} KB  |  Slides: {len(prs.slides)}")
